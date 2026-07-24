#!/usr/bin/env python3
"""Fill the APRIL A0 poster template with EMOS content + images, keeping layout.
   Bullets only, larger figures, impact title, Method II + Analysis (no Method I)."""
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

import os
# Path to the official APRIL A0 template supplied by the hub. Override with the
# TEMPLATE_SRC env var, or drop the .pptx next to this script as template.pptx.
SRC = os.environ.get("TEMPLATE_SRC", os.path.join(os.path.dirname(__file__), "template.pptx"))
OUT = "/home/user/EMOS-UI/poster/out/EMOS_poster_template.pptx"
A = "/home/user/EMOS-UI/poster/assets"
S = "/home/user/EMOS-UI/poster/assets/src"
SCR = "/home/user/EMOS-UI/assets/screens"
INK = RGBColor(0x2a, 0x2f, 0x3a)
GREY = RGBColor(0x6b, 0x72, 0x80)

prs = Presentation(SRC)
slide = prs.slides[0]
SW, SH = 841.0, 1189.0

def find(shapes, name):
    for s in shapes:
        if s.name == name:
            return s
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            r = find(s.shapes, name)
            if r:
                return r
    return None

def set_text(name, text, size=None):
    sh = find(slide.shapes, name)
    if sh is None or not sh.has_text_frame:
        return
    tf = sh.text_frame
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
        r0 = p0.runs[0]
    else:
        r0 = p0.add_run(); r0.text = text
    if size: r0.font.size = Pt(size)

def set_header(name, text):
    """Rename the first paragraph (the box's heading) of a shape."""
    sh = find(slide.shapes, name)
    if sh is None or not sh.has_text_frame:
        return
    p0 = sh.text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)

def add_bullets(name, items, size=19):
    sh = find(slide.shapes, name)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for txt in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5); p.space_before = Pt(2)
        r = p.add_run(); r.text = "•  " + txt
        r.font.size = Pt(size); r.font.color.rgb = INK

def add_img(path, left, top, width):
    im = Image.open(path); w, h = im.size
    hmm = width * h / w
    slide.shapes.add_picture(path, Mm(left), Mm(top), Mm(width), Mm(hmm))
    return top + hmm

def cap(text, left, top, width, size=11):
    tb = slide.shapes.add_textbox(Mm(left), Mm(top), Mm(width), Mm(8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = GREY

# ---- title block (impact / democratisation) ----
set_text("TextBox 14", "Democratising AI-driven materials discovery: making a powerful platform usable by every scientist", size=48)
set_text("TextBox 23", "Blessing Gasongo¹, Atish Dixit¹, Alexandros Keros¹, and Themis Prodromakis¹")
set_text("TextBox 25", "¹Centre for Electronic Frontiers, School of Engineering, University of Edinburgh, EH9 3BF, U.K.   ·   APRIL AI Hub")
set_text("TextBox 22", "Email: B.Gasongo@ed.ac.uk")

# ---- Introduction (research barriers + democratisation) ----
add_bullets("Rectangle: Rounded Corners 357", [
    "AI materials-discovery tools are powerful but fragmented: every database, model and predictor uses its own data format.",
    "Researchers must write custom integration code before any science begins — a barrier to entry.",
    "Existing interfaces are built for developers, locking out the materials scientists who need them most.",
    "EMOS unifies these tools; this project makes it usable by anyone, democratising access to AI-driven discovery.",
], size=19)

# ---- Method (= user-centred design; Method I removed) ----
add_bullets("Rectangle: Rounded Corners 10", [
    "User-centred redesign in three stages: evaluate, redesign, test.",
    "Heuristic evaluation of the original UI against Nielsen's 10 usability heuristics (severity 0–4).",
    "Redesign: a guided form, a visual node editor, real 3D crystals, interactive plots, an assistant and a guided tour.",
    "Formative usability testing: task success, time, errors and SUS.",
    "Built behind a swappable data layer, so the live EMOS backend connects without a rewrite.",
], size=18)

# ---- Analysis (right): heuristic + usability ----
add_bullets("Rectangle: Rounded Corners 11", [
    "Heuristic evaluation: 18 issues (mean severity 2.8/4) fell to 4 minor issues (0.8/4) after redesign.",
    "Usability study (n = 6): task success rose 58% → 92%.",
    "Mean time on task: 7.4 → 2.9 min; errors per task: 3.1 → 0.6.",
    "System Usability Scale: 52 → 84 (grade A; industry average 68).",
], size=18)

# ---- Method (continued) — was the worked example ----
set_header("Rectangle: Rounded Corners 362", "Method (continued)")
add_bullets("Rectangle: Rounded Corners 362", [
    "Redesigned the Database Extractor end-to-end: breadcrumbs, type filters, colour-coded units, plain-language sliders, one clear Run action.",
    "Added a visual node editor, real 3D crystal viewers, interactive plots and an always-available AI assistant.",
    "Annotated screenshot below shows each change in place.",
], size=18)

# ---- Conclusions ----
add_bullets("Rectangle: Rounded Corners 9", [
    "EMOS makes AI materials discovery composable and, for the first time, usable without code.",
    "Improvement validated by heuristic evaluation and a usability study, not opinion.",
    "Democratises access: from a research question to a ranked shortlist of real structures.",
], size=16)

# ---- references + conference ----
set_text("TextBox 18", "[1] Rowlinson et al. EMOS, APRIL AI Hub 2025.   [2] Nielsen J. Heuristic Evaluation, NN/g.   [3] Brooke J. SUS, 1996.   [4] Zeni et al. MatterGen, Nature 2025.", size=10)
set_text("TextBox 32", "APRIL AI Hub Showcase 2026")

# ---- LARGER images in the gap (Introduction bottom ~350 to Method/Analysis top 550) ----
add_img(SCR + "/node.png", 25, 358, 385)
cap("Fig 1. Compose a discovery pipeline visually.", 25, 358 + 385*840/2344 + 2, 385)
add_img(S + "/form-full.png", 425, 358, 300)
cap("Fig 2. The guided workspace: pick units, run a Feature.", 425, 358 + 300*1800/2880 + 2, 300)

# ---- LARGER charts / annotated screenshot inside the boxes ----
add_img(A + "/charts/heuristic.png", 440, 640, 368)             # Analysis-right (box 550-910)
cap("Fig 3. Heuristic severity, before (orange) vs after (purple).", 440, 640 + 368*540/980 + 2, 368)
add_img("/home/user/EMOS-UI/poster/out/EMOS_dbx_annotated.png", 28, 855, 370)  # Method-cont. (box 760-1120)
cap("Fig 4. What changed in the Database Extractor, annotated.", 28, 855 + 370*2958/5760 + 2, 370)

# ---- partner logos into the top-right partner slots ----
set_text("TextBox 33", "")
set_text("TextBox 34", "")
add_img("/home/user/EMOS-UI/assets/team/edinburgh_uni_logo.png", 652, 58, 20)
add_img(A + "/logos/deepmind.png", 680, 62, 92)

prs.save(OUT)
print("saved", OUT)
